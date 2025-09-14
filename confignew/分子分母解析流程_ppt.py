#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

def create_presentation():
    # 创建演示文稿
    prs = Presentation()
    
    # 设置默认字体
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = '微软雅黑'
                    paragraph.font.size = Pt(18)
    
    # 标题页
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "分子分母粒度推断流程"
    subtitle.text = "基于配置文件的自动化推断逻辑"
    
    # 第1页：整体流程概览
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "1. 整体流程概览"
    content = slide.placeholders[1].text_frame
    content.clear()
    
    p = content.add_paragraph()
    p.text = "根据GB列是否存在，分为两种情况："
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "• 无GB列：默认为汇总粒度"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• 有GB列：进一步判断维度汇总或单只"
    p.level = 1
    
    # 第2页：无GB列情况
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "2. 无GB列：汇总粒度"
    content = slide.placeholders[1].text_frame
    content.clear()
    
    p = content.add_paragraph()
    p.text = "当输出列中无GB列时，判定为汇总粒度："
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "• 分子数据项：取 sum:开头的"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• 分母数据项（比例型）：取 mean:开头的"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• 分母粒度：始终为汇总"
    p.level = 1
    
    # 第3页：有GB列 - 维度汇总
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "3. 有GB列：维度汇总判断"
    content = slide.placeholders[1].text_frame
    content.clear()
    
    p = content.add_paragraph()
    p.text = "判断GB列是否属于维度字典："
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "• 若GB列所有值均在维度列表中 → 维度汇总"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• 分子数据项：取 df中 sum:开头的"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• 分母粒度：汇总"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• 分母数据项：取 mean:开头的"
    p.level = 1
    
    # 第4页：有GB列 - 单只判断（比例型）
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "4. 有GB列：单只（比例型）"
    content = slide.placeholders[1].text_frame
    content.clear()
    
    p = content.add_paragraph()
    p.text = "GB列在单只字典中 → 单只（比例型）"
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "• 分子数据项：取 df中 sum:开头的"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• 分母粒度：汇总"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• 分母数据项：取 mean:开头的"
    p.level = 1
    
    # 第5页：有GB列 - 单只判断（非比例型）
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "5. 有GB列：单只（非比例型）"
    content = slide.placeholders[1].text_frame
    content.clear()
    
    p = content.add_paragraph()
    p.text = "GB列在单只字典中 → 单只（非比例型）"
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "根据DF列数量分支判断："
    p.level = 0
    
    p = content.add_paragraph()
    p.text = "• DF只有一项："
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "  - 输出项以 sum:开头 → 数值型"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "  - 输出项以 head:开头 → 禁投型"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "• DF有两项："
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "  - 第一项：head:开头"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "  - 第二项：过滤第一项后剩余的 head:开头"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "• DF有三项："
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "  - 先取 interval:开头的项"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "  - 若含 #（如 A#B）："
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "    - 第一项：#前字段A，head:开头"
    p.level = 3
    
    p = content.add_paragraph()
    p.text = "    - 第二项：#后字段B，head:开头"
    p.level = 3
    
    # 第6页：分子分母相同度量
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "6. 分子分母相同度量"
    content = slide.placeholders[1].text_frame
    content.clear()
    
    p = content.add_paragraph()
    p.text = "当分子分母使用相同度量时："
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "• 特殊情况：比例型但index度量树无out"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "  - 若输出项以 “持仓数量/” 或 “持仓市值/” 开头"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "  - 则分子与分母条件完全相同"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "• 普通情况：根据 filter_expr.expr 中的 out 拆分"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "  - 前半部分：分子条件"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "  - 后半部分：分母条件"
    p.level = 2
    
    # 保存文件到当前目录（确保目录存在）
    output_path = '分子分母解析流程_ppt.pptx'
    prs.save(output_path)
    print(f"✅ PPT 已生成：{output_path}")

if __name__ == "__main__":
    create_presentation()